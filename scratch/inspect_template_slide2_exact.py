import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('Template PPT.pptx')
print(f'Total slides in Template PPT.pptx: {len(prs.slides)}')

for i, slide in enumerate(prs.slides):
    print(f'\n--- Slide {i} ---')
    print(f'  Slide layout: {slide.slide_layout.name}')
    for j, shape in enumerate(slide.shapes):
        text = shape.text_frame.text if shape.has_text_frame else ''
        print(f'  Shape {j} (name={repr(shape.name)}, type={shape.shape_type}): {repr(text[:80])}')
