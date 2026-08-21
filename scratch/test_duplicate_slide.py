import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
import copy

def duplicate_slide_element(prs, source_index=1):
    source_slide = prs.slides[source_index]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Copy background if custom slide background exists
    if hasattr(source_slide, 'background') and source_slide.background:
        try:
            new_slide.background.fill.solid()
            new_slide.background.fill.fore_color.rgb = source_slide.background.fill.fore_color.rgb
        except Exception:
            pass

    return new_slide

prs = Presentation('Template PPT.pptx')
print("Slide 1 layout:", prs.slides[1].slide_layout.name)
print("Slide 1 background:", getattr(prs.slides[1], 'background', None))
