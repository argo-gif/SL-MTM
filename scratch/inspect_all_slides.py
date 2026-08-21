import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('scratch/full_presentation_with_gb_kpi.pptx')
print(f'=== Total Slides Generated: {len(prs.slides)} ===')

for i, slide in enumerate(prs.slides):
    has_bg = slide._element.cSld.bg is not None
    print(f'\n--- Slide {i} (Layout: {slide.slide_layout.name}, HasBG: {has_bg}) ---')
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text_first = shape.text_frame.text.split('\n')[0]
            top = shape.top / 914400.0 if hasattr(shape, 'top') and shape.top is not None else 0
            left = shape.left / 914400.0 if hasattr(shape, 'left') and shape.left is not None else 0
            print(f'  Shape (left={left:.2f}", top={top:.2f}"): {repr(text_first[:100])}')
        elif shape.has_table:
            table = shape.table
            print(f'  Table: rows={len(table.rows)}, cols={len(table.columns)}')
