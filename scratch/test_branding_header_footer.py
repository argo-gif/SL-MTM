import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation('Template PPT.pptx')
slide = prs.slides.add_slide(prs.slide_layouts[2])

# Clear default placeholders
for sp in list(slide.shapes):
    sp._element.getparent().remove(sp._element)

# 1. Top Right Konimex Logo Header
tb_logo = slide.shapes.add_textbox(Inches(7.20), Inches(0.40), Inches(2.20), Inches(0.50))
tf_logo = tb_logo.text_frame
tf_logo.word_wrap = False
p_logo = tf_logo.paragraphs[0]
p_logo.text = "KONIMEX"
p_logo.font.size = Pt(20)
p_logo.font.bold = True
p_logo.font.italic = True
p_logo.font.color.rgb = RGBColor(192, 0, 0)
p_logo.alignment = PP_ALIGN.RIGHT

# 2. Bottom Right Red Footer Banner
banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.50), Inches(4.82), Inches(4.80), Inches(0.38))
banner.fill.solid()
banner.fill.fore_color.rgb = RGBColor(220, 38, 38) # Red Konimex
banner.line.fill.background() # No border line

tf_banner = banner.text_frame
tf_banner.word_wrap = True
tf_banner.margin_top = Pt(3)
tf_banner.margin_bottom = Pt(3)
p_banner = tf_banner.paragraphs[0]
p_banner.text = "Inovasi dan Efisiensi Untuk Pertumbuhan Optimal"
p_banner.font.size = Pt(9.5)
p_banner.font.bold = True
p_banner.font.color.rgb = RGBColor(255, 255, 255) # White
p_banner.alignment = PP_ALIGN.CENTER

prs.save('scratch/test_branding.pptx')
print("Saved scratch/test_branding.pptx successfully!")
