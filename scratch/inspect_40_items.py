import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('scratch/test_40_items_single_slide.pptx')
print(f'=== Total Slides Generated: {len(prs.slides)} ===')

for i in [2, 3]:
    slide = prs.slides[i]
    print(f'\n--- Slide {i} ---')
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text_first = shape.text_frame.text.split('\n')[0]
            top = shape.top / 914400.0 if hasattr(shape, 'top') and shape.top is not None else 0
            left = shape.left / 914400.0 if hasattr(shape, 'left') and shape.left is not None else 0
            print(f'  Shape (left={left:.2f}", top={top:.2f}"): {repr(text_first[:90])}')
        elif shape.has_table:
            table = shape.table
            print(f'  Table: rows={len(table.rows)}, cols={len(table.columns)}')
            if len(table.rows) > 1:
                row1 = [cell.text.replace('\n', ' ') for cell in table.rows[1].cells]
                row_last = [cell.text.replace('\n', ' ') for cell in table.rows[len(table.rows)-1].cells]
                print(f'    Row 1 sample: {row1}')
                print(f'    Row {len(table.rows)-1} sample: {row_last}')
