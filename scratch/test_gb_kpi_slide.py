import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import copy
from backend.data_processor import MTMDataProcessor
from backend.ppt_exporter import format_month_label, build_active_filters_summary

proc = MTMDataProcessor('backend/dataset.db')
filters = {'month': '2026-08', 'mtm_type': 'KA'}
month_label = format_month_label(filters.get('month', '2026-08'))
filter_info = build_active_filters_summary(filters)

gb_grid = proc.get_detail_grid(filters, dimension='grup_brand', limit=200)

prs = Presentation('Template PPT.pptx')
content_layout = prs.slide_layouts[2]

# Insert slide at index 2 (after Slide 1 / Executive KPI Summary)
slide_gb = prs.slides.add_slide(content_layout)

# Clear default placeholder textboxes
for sp in list(slide_gb.shapes):
    sp._element.getparent().remove(sp._element)

# Copy background <p:bg> from Slide 1 if available
if len(prs.slides) > 1 and prs.slides[1]._element.cSld.bg is not None:
    src_bg = prs.slides[1]._element.cSld.bg
    new_bg = copy.deepcopy(src_bg)
    if slide_gb._element.cSld.bg is not None:
        slide_gb._element.cSld.remove(slide_gb._element.cSld.bg)
    slide_gb._element.cSld.insert(0, new_bg)

# Title & Subtitle inside top guide box
title_box = slide_gb.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = f"1.1 SERVICE LEVEL PER GRUP BRAND ({month_label})"
p.font.size = Pt(17)
p.font.bold = True
p.font.color.rgb = RGBColor(192, 0, 0)

p_sub = tf.add_paragraph()
p_sub.text = f"{filter_info} | Performa Pengiriman (SL Kirim) & Realisasi (SL Realisasi) Per Grup Brand"
p_sub.font.size = Pt(8.5)
p_sub.font.bold = True
p_sub.font.color.rgb = RGBColor(100, 100, 100)

# Left Side: Clustered Column Chart (SL Kirim vs SL Realisasi per GB)
chart_data = CategoryChartData()
chart_data.categories = [str(r.get('name', '')) for r in gb_grid]
chart_data.add_series('SL Kirim (%)', [round(float(r.get('sl_kirim', 0)), 1) for r in gb_grid])
chart_data.add_series('SL Realisasi (%)', [round(float(r.get('sl_realisasi', 0)), 1) for r in gb_grid])

cx, cy = Inches(4.50), Inches(3.55)
chart_shape = slide_gb.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.15), cx, cy, chart_data)
chart = chart_shape.chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.TOP
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(8.0)

plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(7.0)
data_labels.font.bold = True

val_axis = chart.value_axis
val_axis.maximum_scale = 100
val_axis.minimum_scale = 0
val_axis.major_unit = 20
val_axis.tick_labels.font.size = Pt(7.5)

cat_axis = chart.category_axis
cat_axis.tick_labels.font.size = Pt(7.5)

if len(chart.series) > 0:
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor(192, 0, 0) # Red Konimex
if len(chart.series) > 1:
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = RGBColor(37, 99, 235) # Royal Blue

# Right Side: Summary Data Table
table_rows = len(gb_grid) + 1
table_cols = 6
table_shape = slide_gb.shapes.add_table(table_rows, table_cols, Inches(5.20), Inches(1.15), Inches(4.25), Inches(3.55))
table = table_shape.table

col_widths = [Inches(0.85), Inches(0.85), Inches(0.85), Inches(0.55), Inches(0.55), Inches(0.60)]
for c_i, w in enumerate(col_widths):
    table.columns[c_i].width = w

headers = ["Grup Brand", "Total Order", "Total Kirim", "SL Kirim", "SL Real.", "GAP"]
for c_i, h in enumerate(headers):
    cell = table.cell(0, c_i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(192, 0, 0)
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(7.5)
        p.alignment = PP_ALIGN.CENTER

for r_i, gb in enumerate(gb_grid):
    r_idx = r_i + 1
    name = str(gb.get('name', ''))
    tp_val = float(gb.get('total_pesan', gb.get('total_p', 0)))
    tk_val = float(gb.get('total_kirim', gb.get('total_k', 0)))
    sl_k = float(gb.get('sl_kirim', 0))
    sl_r = float(gb.get('sl_realisasi', 0))
    gap = sl_r - sl_k

    vals = [
        name,
        f"{tp_val:,.0f}" if tp_val >= 1000 else f"{tp_val:.0f}",
        f"{tk_val:,.0f}" if tk_val >= 1000 else f"{tk_val:.0f}",
        f"{sl_k:.1f}%",
        f"{sl_r:.1f}%",
        f"{gap:+.1f}%"
    ]

    for c_i, val_str in enumerate(vals):
        cell = table.cell(r_idx, c_i)
        cell.text = val_str
        cell.margin_left = Pt(2)
        cell.margin_right = Pt(2)
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(7.0)
            if c_i in [1, 2]:
                p.alignment = PP_ALIGN.RIGHT
            elif c_i in [0, 3, 4, 5]:
                p.alignment = PP_ALIGN.CENTER
            if c_i == 5:
                p.font.bold = True
                p.font.color.rgb = RGBColor(34, 197, 94) if gap >= 0 else RGBColor(220, 38, 38)

prs.save('scratch/test_gb_kpi_slide.pptx')
print("Saved scratch/test_gb_kpi_slide.pptx successfully!")
