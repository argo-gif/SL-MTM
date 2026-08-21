import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from backend.data_processor import MTMDataProcessor

proc = MTMDataProcessor('backend/dataset.db')
filters = {'month': '2026-08', 'mtm_type': 'KA'}

prs = Presentation('Template PPT.pptx')
while len(prs.slides) > 2:
    rId = prs.slides._sldIdLst[2].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[2]

content_layout = prs.slide_layouts[9] if len(prs.slide_layouts) > 9 else prs.slide_layouts[0]

# Test Cabang (16 vital items)
vital_items = proc.get_pareto_tree_maps(filters, dimension='cabang')
vital_items = [item for item in vital_items if (item['cumulative_percentage'] - item['percentage']) < 80.0]

print(f"Cabang vital items count: {len(vital_items)}")

slide_p = prs.slides.add_slide(content_layout)
tb = slide_p.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = f"ANALISIS VITAL PARETO TREEMAP - NAMA CABANG (AGU-2026)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(192, 0, 0)

p_sub = tf.add_paragraph()
p_sub.text = f"📌 Filter Aktif: Bulan = AGU-2026 | MTM = KA | Total Vital: {len(vital_items)} Cabang (80% Unfulfill)"
p_sub.font.size = Pt(8.5)
p_sub.font.bold = True
p_sub.font.color.rgb = RGBColor(100, 100, 100)

cols_cnt = 4
tile_w = Inches(2.05)
tile_h = Inches(0.85)
gap_x = Inches(0.12)
gap_y = Inches(0.08)
top_start = Inches(1.15)

for idx, item in enumerate(vital_items[:16]):
    r_i = idx // cols_cnt
    c_i = idx % cols_cnt

    left_pos = Inches(0.55) + c_i * (tile_w + gap_x)
    top_pos = top_start + r_i * (tile_h + gap_y)

    pct = float(item.get("percentage", 0))
    cum_pct = float(item.get("cumulative_percentage", 0))
    val_num = float(item.get("value", 0))
    val_str = f"Rp {val_num:,.0f}" if val_num >= 1000 else f"{val_num:,.0f} unit"

    shape = slide_p.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, tile_w, tile_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(192, 0, 0) # Red Konimex
    shape.line.color.rgb = RGBColor(234, 179, 8) # Gold border
    shape.line.width = Pt(1.0)

    tf_tile = shape.text_frame
    tf_tile.word_wrap = True
    tf_tile.margin_left = Pt(3)
    tf_tile.margin_right = Pt(3)
    tf_tile.margin_top = Pt(3)
    tf_tile.margin_bottom = Pt(3)

    p0 = tf_tile.paragraphs[0]
    p0.text = f"⭐ VITAL {idx+1}: {item['name']}"
    p0.font.size = Pt(8)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(255, 255, 255)

    p1 = tf_tile.add_paragraph()
    p1.text = val_str
    p1.font.size = Pt(10)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 215, 0)

    p2 = tf_tile.add_paragraph()
    p2.text = f"{pct:.1f}% Kontribusi | {cum_pct:.1f}% Kum."
    p2.font.size = Pt(6.5)
    p2.font.color.rgb = RGBColor(248, 250, 252)

prs.save('scratch/test_16_tiles_grid.pptx')
print('Successfully saved scratch/test_16_tiles_grid.pptx')
