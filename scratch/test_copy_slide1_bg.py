import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
import copy

prs = Presentation('Template PPT.pptx')
source_slide = prs.slides[1]

# Create a new slide with same layout
new_slide = prs.slides.add_slide(source_slide.slide_layout)

# Check if source_slide has background element in XML
s_elem = source_slide._element
bg_elem = s_elem.cSld.bg

if bg_elem is not None:
    print("Found background XML element on source slide!")
    new_bg = copy.deepcopy(bg_elem)
    # Remove existing bg if present, then set
    if new_slide._element.cSld.bg is not None:
        new_slide._element.cSld.remove(new_slide._slide.cSld.bg)
    new_slide._element.cSld.insert(0, new_bg)
    print("Successfully copied background XML element to new slide!")
else:
    print("No bg element on cSld XML, checking relationships...")

prs.save('scratch/test_copied_bg.pptx')
print("Saved scratch/test_copied_bg.pptx!")
