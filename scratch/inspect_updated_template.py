import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('Template PPT.pptx')
print(f'=== Total slides in UPDATED Template PPT.pptx: {len(prs.slides)} ===')

for i, slide in enumerate(prs.slides):
    print(f'\n--- Slide {i} (Layout: {slide.slide_layout.name}) ---')
    print(f'  Background element present: {slide._element.cSld.bg is not None}')
    for j, shape in enumerate(slide.shapes):
        text = shape.text_frame.text if shape.has_text_frame else ''
        top = shape.top / 914400.0 if hasattr(shape, 'top') and shape.top is not None else 0
        left = shape.left / 914400.0 if hasattr(shape, 'left') and shape.left is not None else 0
        w = shape.width / 914400.0 if hasattr(shape, 'width') and shape.width is not None else 0
        h = shape.height / 914400.0 if hasattr(shape, 'height') and shape.height is not None else 0
        print(f'  Shape {j}: name={repr(shape.name)}, type={shape.shape_type}, left={left:.2f}", top={top:.2f}", w={w:.2f}", h={h:.2f}, text={repr(text[:90])}')

print('\n=== Slide Layouts in Template ===')
for i, layout in enumerate(prs.slide_layouts):
    print(f'  Layout {i}: name={layout.name}, shapes={len(layout.shapes)}')
    for shape in layout.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(f'    Layout shape text: {repr(shape.text_frame.text[:80])}')
