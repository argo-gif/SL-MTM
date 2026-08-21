import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from backend.data_processor import MTMDataProcessor
from backend.ppt_exporter import MTMPPTExporter

proc = MTMDataProcessor('backend/dataset.db')
exp = MTMPPTExporter()

# Case 1: ALL Grup Brand
f1 = {'month': '2026-08', 'mtm_type': 'KA'}
data1 = {
    'filters': f1,
    'kpi': proc.get_kpi_scorecard(f1),
    'trend': proc.get_monthly_trend(f1),
    'pareto': {'alasan': proc.get_pareto_tree_maps(f1, dimension='alasan'), 'grup_brand': proc.get_pareto_tree_maps(f1, dimension='grup_brand')},
    'grid_by_dim': {'grup_brand': proc.get_detail_grid(f1, dimension='grup_brand', limit=200)},
    'selected_modules': {'kpi_summary': True, 'pareto_sheets': True, 'detail_grid': True}
}
exp.generate_presentation(data1, 'scratch/test_all_gb_full.pptx')
prs1 = Presentation('scratch/test_all_gb_full.pptx')
gb_titles1 = [s.shapes[0].text_frame.text.split('\n')[0] for s in prs1.slides if s.shapes and s.shapes[0].has_text_frame and "GRUP BRAND" in s.shapes[0].text_frame.text]
print(f"Case 1 (ALL GB): Found {len(gb_titles1)} GRUP BRAND slides -> {gb_titles1}")

# Case 2: Specific Grup Brand ('GB 1')
f2 = {'month': '2026-08', 'mtm_type': 'KA', 'grup_brand': 'GB 1'}
data2 = {
    'filters': f2,
    'kpi': proc.get_kpi_scorecard(f2),
    'trend': proc.get_monthly_trend(f2),
    'pareto': {'alasan': proc.get_pareto_tree_maps(f2, dimension='alasan'), 'grup_brand': proc.get_pareto_tree_maps(f2, dimension='grup_brand')},
    'grid_by_dim': {'grup_brand': proc.get_detail_grid(f2, dimension='grup_brand', limit=200)},
    'selected_modules': {'kpi_summary': True, 'pareto_sheets': True, 'detail_grid': True}
}
exp.generate_presentation(data2, 'scratch/test_specific_gb_full.pptx')
prs2 = Presentation('scratch/test_specific_gb_full.pptx')
gb_titles2 = [s.shapes[0].text_frame.text.split('\n')[0] for s in prs2.slides if s.shapes and s.shapes[0].has_text_frame and "GRUP BRAND" in s.shapes[0].text_frame.text]
print(f"Case 2 (Specific GB='GB 1'): Found {len(gb_titles2)} GRUP BRAND slides -> {gb_titles2}")
