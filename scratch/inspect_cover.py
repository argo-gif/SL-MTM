import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('Template PPT.pptx')
slide0 = prs.slides[0]
print('=== Slide 0 Cover Shapes Inspection ===')
for i, shape in enumerate(slide0.shapes):
    top = shape.top / 914400.0 if hasattr(shape, 'top') and shape.top is not None else 0
    left = shape.left / 914400.0 if hasattr(shape, 'left') and shape.left is not None else 0
    height = shape.height / 914400.0 if hasattr(shape, 'height') and shape.height is not None else 0
    width = shape.width / 914400.0 if hasattr(shape, 'width') and shape.width is not None else 0
    print(f'Shape {i}: name={shape.name}, type={shape.shape_type}, left={left:.2f}", top={top:.2f}", width={width:.2f}", height={height:.2f}"')
    if shape.has_text_frame:
        print('  Text:', repr(shape.text_frame.text[:120]))
