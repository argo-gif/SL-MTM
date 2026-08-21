import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from backend.data_processor import MTMDataProcessor

proc = MTMDataProcessor('backend/dataset.db')
filters = {'month': '2026-08', 'mtm_type': 'KA'}

prs = Presentation('Template PPT.pptx')
# Cleanup extra slides
while len(prs.slides) > 2:
    rId = prs.slides._sldIdLst[2].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[2]

content_layout = prs.slide_layouts[9] if len(prs.slide_layouts) > 9 else prs.slide_layouts[0]

dims = [
    ('alasan', 'ALASAN KETERLAMBATAN'),
    ('mtm_alias', 'AKUN MTM ALIAS'),
    ('cabang', 'NAMA CABANG'),
    ('grup_brand', 'GRUP BRAND'),
    ('item', 'ITEM / PRODUK SKU')
]

for dim_key, dim_label in dims:
    pareto_items = proc.get_pareto_tree_maps(filters, dimension=dim_key)
    # Filter ONLY Vital Pareto 80% items
    vital_items = []
    for item in pareto_items:
        prev_cum = float(item.get('cumulative_percentage', 0)) - float(item.get('percentage', 0))
        if prev_cum < 80.0:
            vital_items.append(item)
            
    if not vital_items and pareto_items:
        vital_items = pareto_items[:1]

    slide_p = prs.slides.add_slide(content_layout)
    
    # Title Box in dashed box at Top: 0.45"
    tb = slide_p.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"ANALISIS VITAL PARETO TREEMAP - {dim_label} (AGU-2026)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)

    p_sub = tf.add_paragraph()
    p_sub.text = f"📌 Filter Aktif: Bulan = AGU-2026 | MTM = KA | Total Vital: {len(vital_items)} Elemen (80% Unfulfill)"
    p_sub.font.size = Pt(8.5)
    p_sub.font.bold = True
    p_sub.font.color.rgb = RGBColor(100, 100, 100)

    n = len(vital_items)
    
    # Determine adaptive grid geometry based on n
    if n <= 3:
        cols_cnt = n
        max_display = n
        tile_w = Inches(8.60 / cols_cnt)
        tile_h = Inches(2.20)
        gap_x = Inches(0.20)
        gap_y = Inches(0.20)
        top_start = Inches(1.35)
        title_font = Pt(11)
        val_font = Pt(15)
        sub_font = Pt(8.5)
    elif n <= 6:
        cols_cnt = 3
        max_display = n
        tile_w = Inches(2.75)
        tile_h = Inches(1.65)
        gap_x = Inches(0.20)
        gap_y = Inches(0.18)
        top_start = Inches(1.15)
        title_font = Pt(10)
        val_font = Pt(13)
        sub_font = Pt(7.5)
    elif n <= 9:
        cols_cnt = 3
        max_display = n
        tile_w = Inches(2.75)
        tile_h = Inches(1.12)
        gap_x = Inches(0.20)
        gap_y = Inches(0.12)
        top_start = Inches(1.15)
        title_font = Pt(9)
        val_font = Pt(11.5)
        sub_font = Pt(7)
    else: # n >= 10
        cols_cnt = 4
        max_display = min(12, n)
        tile_w = Inches(2.05)
        tile_h = Inches(1.10)
        gap_x = Inches(0.15)
        gap_y = Inches(0.12)
        top_start = Inches(1.15)
        title_font = Pt(8.5)
        val_font = Pt(10.5)
        sub_font = Pt(6.5)

    display_items = vital_items[:max_display]
    print(f"Dimensi {dim_label}: {n} vital items, displaying {len(display_items)} tiles (cols={cols_cnt})")

    for idx, item in enumerate(display_items):
        r_i = idx // cols_cnt
        c_i = idx % cols_cnt

        left_pos = Inches(0.55) + c_i * (tile_w + gap_x)
        top_pos = top_start + r_i * (tile_h + gap_y)

        pct = float(item.get("percentage", 0))
        cum_pct = float(item.get("cumulative_percentage", 0))
        val_num = float(item.get("value", 0))
        val_str = f"Rp {val_num:,.0f}" if val_num >= 1000 else f"{val_num:,.0f} unit"

        shape = slide_p.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, tile_w, tile_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(192, 0, 0) # Red Konimex
        shape.line.color.rgb = RGBColor(234, 179, 8) # Gold border
        shape.line.width = Pt(1.2)

        tf_tile = shape.text_frame
        tf_tile.word_wrap = True
        tf_tile.margin_left = Pt(4)
        tf_tile.margin_right = Pt(4)
        tf_tile.margin_top = Pt(4)
        tf_tile.margin_bottom = Pt(4)

        # Line 0: Badge Status
        p0 = tf_tile.paragraphs[0]
        p0.text = "⭐ VITAL PARETO 80%"
        p0.font.size = sub_font
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(255, 215, 0)

        # Line 1: Name
        p1 = tf_tile.add_paragraph()
        p1.text = str(item.get("name", "-"))
        p1.font.size = title_font
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)

        # Line 2: Value
        p2 = tf_tile.add_paragraph()
        p2.text = val_str
        p2.font.size = val_font
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)

        # Line 3: Pct & Cum Pct
        p3 = tf_tile.add_paragraph()
        p3.text = f"{pct:.1f}% Kontribusi | {cum_pct:.1f}% Kum."
        p3.font.size = sub_font
        p3.font.color.rgb = RGBColor(248, 250, 252)

prs.save('scratch/test_adaptive_treemap.pptx')
print('Successfully saved scratch/test_adaptive_treemap.pptx')
