import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('Template PPT.pptx')
s2 = prs.slides[2]
print('Slide 2 layout:', s2.slide_layout.name)
print('Slide 2 shapes count:', len(s2.shapes))
for i, sh in enumerate(s2.shapes):
    txt = sh.text_frame.text if sh.has_text_frame else ''
    print(f'Shape {i}: type={sh.shape_type}, text={repr(txt)}')
