import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation('Template PPT.pptx')
print(f"Total slides in updated template: {len(prs.slides)}")

# Inspect Slide 1 and Slide 2 XML
for idx in [0, 1, 2, 14]:
    s = prs.slides[idx]
    bg = s._element.cSld.bg
    print(f"Slide {idx} bg tag: {bg.tag if bg is not None else None}")
