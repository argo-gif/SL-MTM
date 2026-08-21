import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('scratch/aligned_dashed_box.pptx')
slide = prs.slides[1]
print('=== Slide 1 Aligned Dashed Box Inspection ===')
for i, shape in enumerate(slide.shapes):
    top = shape.top / 914400.0 if hasattr(shape, 'top') and shape.top is not None else 0
    left = shape.left / 914400.0 if hasattr(shape, 'left') and shape.left is not None else 0
    height = shape.height / 914400.0 if hasattr(shape, 'height') and shape.height is not None else 0
    width = shape.width / 914400.0 if hasattr(shape, 'width') and shape.width is not None else 0
    print(f'Shape {i}: type={shape.shape_type}, left={left:.2f}", top={top:.2f}", width={width:.2f}", height={height:.2f}"')
    if shape.has_text_frame:
        print('  Text:', repr(shape.text_frame.text[:120]))
    if shape.has_chart:
        chart = shape.chart
        print(f'  Chart type={chart.chart_type}, series_cnt={len(chart.series)}')
        for s_idx, s in enumerate(chart.series):
            print(f'    Series {s_idx}: name={s.name}, values={s.values}')
