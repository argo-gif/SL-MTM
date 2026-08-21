import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('Template PPT.pptx')
s1 = prs.slides[1]
print(f'Slide 1 shapes count: {len(s1.shapes)}')

for i, shape in enumerate(s1.shapes):
    text = shape.text_frame.text if shape.has_text_frame else ''
    top = shape.top / 914400.0 if hasattr(shape, 'top') and shape.top is not None else 0
    left = shape.left / 914400.0 if hasattr(shape, 'left') and shape.left is not None else 0
    width = shape.width / 914400.0 if hasattr(shape, 'width') and shape.width is not None else 0
    height = shape.height / 914400.0 if hasattr(shape, 'height') and shape.height is not None else 0
    print(f'Shape {i}: name={repr(shape.name)}, type={shape.shape_type}, left={left:.2f}", top={top:.2f}", w={width:.2f}", h={height:.2f}", text={repr(text)}')
