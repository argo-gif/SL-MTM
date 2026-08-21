import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
import copy

def clone_slide(prs, source_slide_index=1):
    source_slide = prs.slides[source_slide_index]
    # Create a new blank slide or slide with same layout
    blank_layout = prs.slide_layouts[2]
    new_slide = prs.slides.add_slide(blank_layout)

    # Copy shapes from source_slide to new_slide XML elements
    # Or copy shapes from source slide element
    for shape in source_slide.shapes:
        # We can copy elements or copy shape XML
        new_sp = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_sp, 'p:extLst')

    return new_slide

prs = Presentation('Template PPT.pptx')
print(f'Initial slides: {len(prs.slides)}')
new_s = clone_slide(prs, 1)
print(f'Slides after clone: {len(prs.slides)}')
print(f'Cloned slide shapes count: {len(new_s.shapes)}')
prs.save('scratch/test_cloned_slide.pptx')
print('Saved scratch/test_cloned_slide.pptx!')
