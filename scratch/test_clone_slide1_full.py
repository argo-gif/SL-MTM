import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
import copy

def create_slide_cloned_from_slide1(prs):
    source_slide = prs.slides[1]
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    
    # 1. Clear any default shapes from new_slide
    for sp in list(new_slide.shapes):
        sp._element.getparent().remove(sp._element)

    # 2. Copy background <p:bg> element if present on source_slide
    src_cSld = source_slide._element.cSld
    if src_cSld.bg is not None:
        new_bg = copy.deepcopy(src_cSld.bg)
        if new_slide._element.cSld.bg is not None:
            new_slide._element.cSld.remove(new_slide._element.cSld.bg)
        new_slide._element.cSld.insert(0, new_bg)

    # 3. Copy any shapes from source_slide if any exist
    for shape in source_slide.shapes:
        if shape.name != 'Title 4': # Keep title clear for new slide title
            new_sp = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.append(new_sp)

    return new_slide

prs = Presentation('Template PPT.pptx')
s_new = create_slide_cloned_from_slide1(prs)
print("Created cloned slide successfully!")
prs.save('scratch/test_full_cloned_slide1.pptx')
print("Saved scratch/test_full_cloned_slide1.pptx!")
