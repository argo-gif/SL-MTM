import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation('Template PPT.pptx')
slide = prs.slides.add_slide(prs.slide_layouts[9])

# Title Box in dashed box at Top: 0.45"
tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '2.1 ANALISIS PARETO TREEMAP - ALASAN KETERLAMBATAN (AGU-2026)'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(192, 0, 0)

p_sub = tf.add_paragraph()
p_sub.text = '📌 Filter Aktif: Bulan = AGU-2026 | MTM = KA'
p_sub.font.size = Pt(8.5)
p_sub.font.bold = True
p_sub.font.color.rgb = RGBColor(100, 100, 100)

pareto_items = [
    {'name': 'STOCK OUT PRINCIPAL', 'value': 699049730, 'percentage': 36.1, 'cumulative_percentage': 36.1},
    {'name': 'STOCK OUT CABANG', 'value': 353910900, 'percentage': 18.3, 'cumulative_percentage': 54.4},
    {'name': 'BARANG DALAM PERJALANAN (BDP)', 'value': 234786730, 'percentage': 12.1, 'cumulative_percentage': 66.5},
    {'name': 'PO SALAH HARGA/DISC', 'value': 214066690, 'percentage': 11.1, 'cumulative_percentage': 77.6},
    {'name': 'PERMINTAAN PLG TERIMA SEBAGIAN', 'value': 136485630, 'percentage': 7.0, 'cumulative_percentage': 84.6},
    {'name': 'PRODUK SUDAH TIDAK PRODUKSI', 'value': 90798820, 'percentage': 4.7, 'cumulative_percentage': 89.3}
]

# Render Treemap Tile Layout
cols = 3
tile_w = Inches(2.80)
tile_h = Inches(1.80)
gap_x = Inches(0.20)
gap_y = Inches(0.20)

for idx, item in enumerate(pareto_items[:6]):
    r = idx // cols
    c = idx % cols
    
    left = Inches(0.55) + c * (tile_w + gap_x)
    top = Inches(1.20) + r * (tile_h + gap_y)
    
    cum_pct = item['cumulative_percentage']
    pct = item['percentage']
    prev_cum = cum_pct - pct
    is_vital = (prev_cum < 80.0)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, tile_w, tile_h)
    shape.fill.solid()
    if is_vital:
        shape.fill.fore_color.rgb = RGBColor(192, 0, 0) # Red Konimex
        shape.line.color.rgb = RGBColor(234, 179, 8) # Gold border
        shape.line.width = Pt(1.5)
    else:
        shape.fill.fore_color.rgb = RGBColor(100, 116, 139) # Slate Muted
        shape.line.color.rgb = RGBColor(148, 163, 184)
    
    tf_tile = shape.text_frame
    tf_tile.word_wrap = True
    
    # Line 0: Badge Status
    p0 = tf_tile.paragraphs[0]
    p0.text = '⭐ VITAL PARETO 80%' if is_vital else 'TRIVIAL 20%'
    p0.font.size = Pt(8)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(255, 215, 0) if is_vital else RGBColor(226, 232, 240)
    
    # Line 1: Name
    p1 = tf_tile.add_paragraph()
    p1.text = item['name']
    p1.font.size = Pt(10.5)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    
    # Line 2: Value
    p2 = tf_tile.add_paragraph()
    p2.text = f"Rp {item['value']:,.0f}"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    
    # Line 3: Pct & Cum Pct
    p3 = tf_tile.add_paragraph()
    p3.text = f"{pct:.1f}% Kontribusi | {cum_pct:.1f}% Kumulatif"
    p3.font.size = Pt(8)
    p3.font.color.rgb = RGBColor(248, 250, 252)

prs.save('scratch/test_treemap_ppt.pptx')
print('Successfully generated scratch/test_treemap_ppt.pptx')
