import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from backend.data_processor import MTMDataProcessor

proc = MTMDataProcessor('backend/dataset.db')
filters = {'month': '2026-08', 'mtm_type': 'KA'}

prs = Presentation('Template PPT.pptx')
while len(prs.slides) > 2:
    rId = prs.slides._sldIdLst[2].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[2]

content_layout = prs.slide_layouts[9] if len(prs.slide_layouts) > 9 else prs.slide_layouts[0]

vital_items = proc.get_pareto_tree_maps(filters, dimension='item')
vital_items = [item for item in vital_items if (item['cumulative_percentage'] - item['percentage']) < 80.0]

print(f"Total Item SKU Vital Pareto items: {len(vital_items)}")

# 1. Slide 6.1: Treemap Slide (All 40 Vital Items on 1 Slide)
slide_p = prs.slides.add_slide(content_layout)
tb = slide_p.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = f"6.1 ANALISIS PARETO UNFULLFILL - ITEM / PRODUK SKU (AGU-2026)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(192, 0, 0)

p_sub = tf.add_paragraph()
p_sub.text = f"📌 Filter Aktif: Bulan = AGU-2026 | MTM = KA | Total Vital Pareto: {len(vital_items)} Item SKU (80% Unfulfill)"
p_sub.font.size = Pt(8.5)
p_sub.font.bold = True
p_sub.font.color.rgb = RGBColor(100, 100, 100)

cols_cnt = 5
tile_w = Inches(1.64)
tile_h = Inches(0.43)
gap_x = Inches(0.10)
gap_y = Inches(0.04)
top_start = Inches(1.15)

for idx, item in enumerate(vital_items):
    r_i = idx // cols_cnt
    c_i = idx % cols_cnt

    left_pos = Inches(0.55) + c_i * (tile_w + gap_x)
    top_pos = top_start + r_i * (tile_h + gap_y)

    pct = float(item.get("percentage", 0))
    cum_pct = float(item.get("cumulative_percentage", 0))
    val_num = float(item.get("value", 0))
    val_str = f"Rp {val_num:,.0f}" if val_num >= 1000 else f"{val_num:,.0f}"

    shape = slide_p.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, tile_w, tile_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(192, 0, 0) # Red Konimex
    shape.line.color.rgb = RGBColor(234, 179, 8) # Gold border
    shape.line.width = Pt(0.8)

    tf_tile = shape.text_frame
    tf_tile.word_wrap = True
    tf_tile.margin_left = Pt(2)
    tf_tile.margin_right = Pt(2)
    tf_tile.margin_top = Pt(1)
    tf_tile.margin_bottom = Pt(1)

    p0 = tf_tile.paragraphs[0]
    p0.text = f"#{idx+1} {str(item.get('name', '-'))}"
    p0.font.size = Pt(6.5)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(255, 255, 255)

    p1 = tf_tile.add_paragraph()
    p1.text = f"{val_str} ({pct:.1f}%)"
    p1.font.size = Pt(6.0)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 215, 0)

# 2. Slide 6.2: Detail Data Grid Slide (All 40 Vital Rows on 1 Slide)
grid_dim = proc.get_detail_grid(filters, dimension='item', limit=200)
vital_names = set([v.get("name") for v in vital_items if v.get("name")])
vital_grid = [g for g in grid_dim if g.get("name") in vital_names]

slide_g = prs.slides.add_slide(content_layout)
title_box_g = slide_g.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
tf_g = title_box_g.text_frame
tf_g.word_wrap = True
p_g = tf_g.paragraphs[0]
p_g.text = f"6.2 TABEL DETAIL TRANSAKSI - VITAL PARETO ITEM / PRODUK SKU (AGU-2026)"
p_g.font.size = Pt(15)
p_g.font.bold = True
p_g.font.color.rgb = RGBColor(192, 0, 0)

p_sub_g = tf_g.add_paragraph()
p_sub_g.text = f"📌 Filter Aktif: Bulan = AGU-2026 | MTM = KA | Filter Vital Pareto 80% ({len(vital_grid)} Item)"
p_sub_g.font.size = Pt(8.5)
p_sub_g.font.bold = True
p_sub_g.font.color.rgb = RGBColor(100, 100, 100)

rows_g_cnt = len(vital_grid) + 1
g_table = slide_g.shapes.add_table(rows_g_cnt, 8, Inches(0.55), Inches(1.15), Inches(8.8), Inches(3.8)).table

g_col_widths = [Inches(0.4), Inches(2.7), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.1), Inches(0.8), Inches(0.8)]
for c_i, w in enumerate(g_col_widths):
    g_table.columns[c_i].width = w

headers_g = ["#", "Nama Item SKU (Vital 80%)", "Total Order", "Total Kirim", "Total Realisasi", "Nilai Unfulfill", "SL Kirim", "SL Realisasi"]
for c, h in enumerate(headers_g):
    cell = g_table.cell(0, c)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(192, 0, 0)
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(6.5)
        p.alignment = PP_ALIGN.CENTER

for r, row_data in enumerate(vital_grid):
    r_idx = r + 1
    sl_k_val = float(row_data.get("sl_kirim", 0))
    sl_r_val = float(row_data.get("sl_realisasi", 0))
    
    tp_val = row_data.get('total_p', row_data.get('total_pesan', 0))
    tk_val = row_data.get('total_k', row_data.get('total_kirim', 0))
    tr_val = row_data.get('total_r', row_data.get('total_realisasi', 0))
    gap_val = row_data.get('gap_unfulfill', 0)

    vals = [
        str(r_idx),
        str(row_data.get("name", "")),
        f"{tp_val:,.0f}",
        f"{tk_val:,.0f}",
        f"{tr_val:,.0f}",
        f"{gap_val:,.0f}",
        f"{sl_k_val:.1f}%",
        f"{sl_r_val:.1f}%"
    ]
    for c, v in enumerate(vals):
        cell = g_table.cell(r_idx, c)
        cell.text = v
        cell.margin_left = Pt(2)
        cell.margin_right = Pt(2)
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(6.0)
            if c in [2, 3, 4, 5]:
                p.alignment = PP_ALIGN.RIGHT
            elif c in [0, 6, 7]:
                p.alignment = PP_ALIGN.CENTER

prs.save('scratch/test_40_items_single_slide.pptx')
print('Successfully saved scratch/test_40_items_single_slide.pptx')
