import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('Template PPT.pptx')
s1 = prs.slides[1]
print(f'Slide 1 shapes count: {len(s1.shapes)}')
for i, sh in enumerate(s1.shapes):
    text = sh.text_frame.text if sh.has_text_frame else ''
    print(f'Shape {i}: name={repr(sh.name)}, type={sh.shape_type}, text={repr(text)}')
