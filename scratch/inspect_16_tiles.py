import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('scratch/test_16_tiles_grid.pptx')
slide = prs.slides[2]
print('=== 16 Tiles Slide Inspection ===')
for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame and shape.text_frame.text.strip():
        top = shape.top / 914400.0 if hasattr(shape, 'top') and shape.top is not None else 0
        left = shape.left / 914400.0 if hasattr(shape, 'left') and shape.left is not None else 0
        print(f'  Shape {i} (left={left:.2f}", top={top:.2f}"): {repr(shape.text_frame.text.splitlines()[0])}')
