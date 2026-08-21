import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from backend.data_processor import MTMDataProcessor
from backend.ppt_exporter import MTMPPTExporter

proc = MTMDataProcessor('backend/dataset.db')
exp = MTMPPTExporter()

# Case 1: ALL Grup Brand selected
f1 = {'month': '2026-08', 'mtm_type': 'KA'}
data1 = {
    'filters': f1,
    'kpi': proc.get_kpi_scorecard(f1),
    'trend': proc.get_monthly_trend(f1),
    'pareto': {'alasan': proc.get_pareto_tree_maps(f1, dimension='alasan')},
    'grid_by_dim': {'grup_brand': proc.get_detail_grid(f1, dimension='grup_brand', limit=200)},
    'selected_modules': {'kpi_summary': True, 'pareto_sheets': True, 'detail_grid': True}
}
exp.generate_presentation(data1, 'scratch/test_all_gb.pptx')
prs1 = Presentation('scratch/test_all_gb.pptx')
titles1 = [s.shapes[0].text_frame.text.split('\n')[0] for s in prs1.slides if s.shapes and s.shapes[0].has_text_frame]
has_gb_slide1 = any("1.1 SERVICE LEVEL PER GRUP BRAND" in t for t in titles1)
print(f"Case 1 (ALL GB): Has Slide 1.1 = {has_gb_slide1} (Total slides: {len(prs1.slides)})")

# Case 2: Specific Grup Brand selected ('GB 1')
f2 = {'month': '2026-08', 'mtm_type': 'KA', 'grup_brand': 'GB 1'}
data2 = {
    'filters': f2,
    'kpi': proc.get_kpi_scorecard(f2),
    'trend': proc.get_monthly_trend(f2),
    'pareto': {'alasan': proc.get_pareto_tree_maps(f2, dimension='alasan')},
    'grid_by_dim': {'grup_brand': proc.get_detail_grid(f2, dimension='grup_brand', limit=200)},
    'selected_modules': {'kpi_summary': True, 'pareto_sheets': True, 'detail_grid': True}
}
exp.generate_presentation(data2, 'scratch/test_specific_gb.pptx')
prs2 = Presentation('scratch/test_specific_gb.pptx')
titles2 = [s.shapes[0].text_frame.text.split('\n')[0] for s in prs2.slides if s.shapes and s.shapes[0].has_text_frame]
has_gb_slide2 = any("1.1 SERVICE LEVEL PER GRUP BRAND" in t for t in titles2)
print(f"Case 2 (Specific GB='GB 1'): Has Slide 1.1 = {has_gb_slide2} (Total slides: {len(prs2.slides)})")
