import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from backend.data_processor import MTMDataProcessor

proc = MTMDataProcessor('backend/dataset.db')
filters = {'month': '2026-08', 'mtm_type': 'KA'}

prs = Presentation('Template PPT.pptx')
while len(prs.slides) > 2:
    rId = prs.slides._sldIdLst[2].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[2]

content_layout = prs.slide_layouts[9] if len(prs.slide_layouts) > 9 else prs.slide_layouts[0]

vital_items = proc.get_pareto_tree_maps(filters, dimension='item')
vital_items = [item for item in vital_items if (item['cumulative_percentage'] - item['percentage']) < 80.0]

print(f"Total Item SKU Vital Pareto items: {len(vital_items)}")

# Chunk size for large dimensions
CHUNK_SIZE = 20
chunks = [vital_items[i:i + CHUNK_SIZE] for i in range(0, len(vital_items), CHUNK_SIZE)]
total_chunks = len(chunks)

print(f"Total Treemap chunks generated: {total_chunks}")

# Render Treemap Chunks
for chunk_idx, chunk_items in enumerate(chunks):
    slide_p = prs.slides.add_slide(content_layout)
    tb = slide_p.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    part_suffix = f" (BAGIAN {chunk_idx+1}/{total_chunks})" if total_chunks > 1 else ""
    p.text = f"6.1 ANALISIS PARETO UNFULLFILL - ITEM / PRODUK SKU{part_suffix} (AGU-2026)"
    p.font.size = Pt(15 if total_chunks > 1 else 16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)

    p_sub = tf.add_paragraph()
    start_num = chunk_idx * CHUNK_SIZE + 1
    end_num = start_num + len(chunk_items) - 1
    p_sub.text = f"📌 Filter Aktif: Bulan = AGU-2026 | MTM = KA | Item Vital #{start_num} - #{end_num} (Total Vital: {len(vital_items)})"
    p_sub.font.size = Pt(8.5)
    p_sub.font.bold = True
    p_sub.font.color.rgb = RGBColor(100, 100, 100)

    cols_cnt = 4
    tile_w = Inches(2.05)
    tile_h = Inches(0.68)
    gap_x = Inches(0.12)
    gap_y = Inches(0.06)
    top_start = Inches(1.15)

    for idx, item in enumerate(chunk_items):
        item_global_idx = start_num + idx
        r_i = idx // cols_cnt
        c_i = idx % cols_cnt

        left_pos = Inches(0.55) + c_i * (tile_w + gap_x)
        top_pos = top_start + r_i * (tile_h + gap_y)

        pct = float(item.get("percentage", 0))
        cum_pct = float(item.get("cumulative_percentage", 0))
        val_num = float(item.get("value", 0))
        val_str = f"Rp {val_num:,.0f}" if val_num >= 1000 else f"{val_num:,.0f}"

        shape = slide_p.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, tile_w, tile_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(192, 0, 0) # Red Konimex
        shape.line.color.rgb = RGBColor(234, 179, 8) # Gold border
        shape.line.width = Pt(1.0)

        tf_tile = shape.text_frame
        tf_tile.word_wrap = True
        tf_tile.margin_left = Pt(3)
        tf_tile.margin_right = Pt(3)
        tf_tile.margin_top = Pt(2)
        tf_tile.margin_bottom = Pt(2)

        p0 = tf_tile.paragraphs[0]
        p0.text = f"⭐ VITAL #{item_global_idx}: {str(item.get('name', '-'))}"
        p0.font.size = Pt(7.5)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(255, 255, 255)

        p1 = tf_tile.add_paragraph()
        p1.text = val_str
        p1.font.size = Pt(9.5)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 215, 0)

        p2 = tf_tile.add_paragraph()
        p2.text = f"{pct:.1f}% Kontribusi | {cum_pct:.1f}% Kum."
        p2.font.size = Pt(6.5)
        p2.font.color.rgb = RGBColor(248, 250, 252)

# Detail Data Grid Chunks
grid_dim = proc.get_detail_grid(filters, dimension='item', limit=200)
vital_names = set([v.get("name") for v in vital_items if v.get("name")])
vital_grid = [g for g in grid_dim if g.get("name") in vital_names]

grid_chunks = [vital_grid[i:i + CHUNK_SIZE] for i in range(0, len(vital_grid), CHUNK_SIZE)]
total_grid_chunks = len(grid_chunks)

for chunk_idx, chunk_grid in enumerate(grid_chunks):
    slide_g = prs.slides.add_slide(content_layout)
    title_box_g = slide_g.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
    tf_g = title_box_g.text_frame
    tf_g.word_wrap = True
    p_g = tf_g.paragraphs[0]
    part_suffix = f" (BAGIAN {chunk_idx+1}/{total_grid_chunks})" if total_grid_chunks > 1 else ""
    p_g.text = f"6.2 TABEL DETAIL TRANSAKSI - VITAL PARETO ITEM / PRODUK SKU{part_suffix} (AGU-2026)"
    p_g.font.size = Pt(14 if total_grid_chunks > 1 else 15)
    p_g.font.bold = True
    p_g.font.color.rgb = RGBColor(192, 0, 0)

    p_sub_g = tf_g.add_paragraph()
    start_num = chunk_idx * CHUNK_SIZE + 1
    end_num = start_num + len(chunk_grid) - 1
    p_sub_g.text = f"📌 Filter Aktif: Bulan = AGU-2026 | MTM = KA | Filter Vital Pareto 80% (#{start_num} - #{end_num})"
    p_sub_g.font.size = Pt(8.5)
    p_sub_g.font.bold = True
    p_sub_g.font.color.rgb = RGBColor(100, 100, 100)

    rows_g_cnt = len(chunk_grid) + 1
    g_table = slide_g.shapes.add_table(rows_g_cnt, 8, Inches(0.55), Inches(1.15), Inches(8.8), Inches(3.8)).table

    g_col_widths = [Inches(0.5), Inches(2.3), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1), Inches(0.8), Inches(0.8)]
    for c_i, w in enumerate(g_col_widths):
        g_table.columns[c_i].width = w

    headers_g = ["#", "Nama Item SKU (Vital 80%)", "Total Order", "Total Kirim", "Total Realisasi", "Nilai Unfulfill", "SL Kirim", "SL Realisasi"]
    for c, h in enumerate(headers_g):
        cell = g_table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(192, 0, 0)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(8.0)
            p.alignment = PP_ALIGN.CENTER

    for r, row_data in enumerate(chunk_grid):
        r_idx = r + 1
        global_r_idx = start_num + r
        sl_k_val = float(row_data.get("sl_kirim", 0))
        sl_r_val = float(row_data.get("sl_realisasi", 0))
        
        tp_val = row_data.get('total_p', row_data.get('total_pesan', 0))
        tk_val = row_data.get('total_k', row_data.get('total_kirim', 0))
        tr_val = row_data.get('total_r', row_data.get('total_realisasi', 0))
        gap_val = row_data.get('gap_unfulfill', 0)

        vals = [
            str(global_r_idx),
            str(row_data.get("name", "")),
            f"{tp_val:,.0f}",
            f"{tk_val:,.0f}",
            f"{tr_val:,.0f}",
            f"{gap_val:,.0f}",
            f"{sl_k_val:.1f}%",
            f"{sl_r_val:.1f}%"
        ]
        for c, v in enumerate(vals):
            cell = g_table.cell(r_idx, c)
            cell.text = v
            cell.margin_left = Pt(3)
            cell.margin_right = Pt(3)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7.5)
                if c in [2, 3, 4, 5]:
                    p.alignment = PP_ALIGN.RIGHT
                elif c in [0, 6, 7]:
                    p.alignment = PP_ALIGN.CENTER

prs.save('scratch/test_split_2_slides.pptx')
print('Successfully saved scratch/test_split_2_slides.pptx')
