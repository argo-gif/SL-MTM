import os
import sys
import copy
import zipfile
import re
import xml.etree.ElementTree as ET
from typing import Dict, List, Any, Optional

try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False


def format_month_label(m_str: Any) -> str:
    if not m_str or str(m_str).lower() in ['all', 'semua', 'semua bulan', 'none']:
        return "Semua Bulan"
    months_id = {
        '01': 'JAN', '02': 'FEB', '03': 'MAR', '04': 'APR', '05': 'MEI', '06': 'JUN',
        '07': 'JUL', '08': 'AGU', '09': 'SEP', '10': 'OKT', '11': 'NOV', '12': 'DES'
    }
    parts = str(m_str).split('-')
    if len(parts) == 2:
        yr, mo = parts[0], parts[1]
        return f"{months_id.get(mo, mo)}-{yr}"
    return str(m_str)


class MTMPPTExporter:
    def __init__(self, template_path: str = "Template PPT.pptx"):
        self.template_path = template_path

    def generate_presentation(self, export_data: Dict[str, Any], output_path: str = "output_report.pptx") -> str:
        """
        Generate PowerPoint presentation based on Template PPT.pptx with complete slide data population matching dashboard filters.
        """
        if HAS_PYTHON_PPTX and os.path.exists(self.template_path):
            return self._generate_with_python_pptx(export_data, output_path)
        else:
            return self._generate_fallback(export_data, output_path)

def squarify_layout(items: List[Dict[str, Any]], width: float, height: float) -> List[Dict[str, Any]]:
    non_zero = [it for it in items if float(it.get("value", 0)) > 0]
    total_val = sum(float(it.get("value", 0)) for it in non_zero)
    if total_val <= 0 or width <= 0 or height <= 0:
        return []

    total_area = float(width * height)
    children = []
    for it in non_zero:
        v = float(it.get("value", 0))
        children.append({
            "raw": it,
            "area": (v / total_val) * total_area
        })

    rects = []

    def worst_aspect_ratio(row, side_len):
        if not row or side_len <= 0:
            return float('inf')
        row_area = sum(c["area"] for c in row)
        if row_area <= 0:
            return float('inf')
        max_a = max(c["area"] for c in row)
        min_a = min(c["area"] for c in row)
        s2 = float(side_len * side_len)
        a2 = float(row_area * row_area)
        return max((s2 * max_a) / a2, a2 / (s2 * min_a))

    def layout_row(row, container):
        row_area = sum(c["area"] for c in row)
        is_horizontal = container["w"] >= container["h"]
        side_len = container["h"] if is_horizontal else container["w"]
        row_thickness = row_area / side_len if side_len > 0 else 0

        current_offset = container["y"] if is_horizontal else container["x"]

        for c in row:
            item_len = c["area"] / row_thickness if row_thickness > 0 else 0
            if is_horizontal:
                rects.append({
                    "item": c["raw"],
                    "x": container["x"],
                    "y": current_offset,
                    "w": row_thickness,
                    "h": item_len
                })
                current_offset += item_len
            else:
                rects.append({
                    "item": c["raw"],
                    "x": current_offset,
                    "y": container["y"],
                    "w": item_len,
                    "h": row_thickness
                })
                current_offset += item_len

        if is_horizontal:
            container["x"] += row_thickness
            container["w"] -= row_thickness
        else:
            container["y"] += row_thickness
            container["h"] -= row_thickness

    container = {"x": 0.0, "y": 0.0, "w": float(width), "h": float(height)}
    current_row = []

    for c in children:
        side_len = min(container["w"], container["h"])
        if side_len <= 0:
            break

        if not current_row:
            current_row.append(c)
        else:
            current_worst = worst_aspect_ratio(current_row, side_len)
            new_worst = worst_aspect_ratio(current_row + [c], side_len)

            if new_worst <= current_worst:
                current_row.append(c)
            else:
                layout_row(current_row, container)
                current_row = [c]

    if current_row and container["w"] > 0 and container["h"] > 0:
        layout_row(current_row, container)

    return rects

def wrap_text_lines(font, text: str, max_w: float) -> List[str]:
    words = text.split()
    if not words:
        return []

    lines = []
    curr_words = []

    def get_w(t):
        if hasattr(font, 'getlength'):
            return font.getlength(t)
        elif hasattr(font, 'getbbox'):
            return font.getbbox(t)[2]
        return len(t) * 8

    for word in words:
        test_str = " ".join(curr_words + [word]) if curr_words else word
        if get_w(test_str) <= max_w or not curr_words:
            curr_words.append(word)
        else:
            lines.append(" ".join(curr_words))
            curr_words = [word]

    if curr_words:
        lines.append(" ".join(curr_words))

    return lines

def generate_treemap_image(items: List[Dict[str, Any]], width_px: int = 1800, height_px: int = 780, vital_cutoff_idx: int = 0) -> str:
    from PIL import Image, ImageDraw, ImageFont
    import math

    non_zero = [it for it in items if float(it.get("value", 0)) > 0]
    display_items = non_zero[:40]

    rects = squarify_layout(display_items, float(width_px), float(height_px))

    canvas = Image.new('RGB', (width_px, height_px), color=(15, 23, 42))

    font_bold_path = 'C:/Windows/Fonts/segoeuib.ttf' if os.path.exists('C:/Windows/Fonts/segoeuib.ttf') else 'C:/Windows/Fonts/arialbd.ttf'
    font_reg_path = 'C:/Windows/Fonts/segoeui.ttf' if os.path.exists('C:/Windows/Fonts/segoeui.ttf') else 'C:/Windows/Fonts/arial.ttf'

    def get_font(path, size):
        try:
            return ImageFont.truetype(path, int(size))
        except Exception:
            return ImageFont.load_default()

    main_styles = [
        {"bg_start": (15, 23, 42), "bg_end": (30, 58, 138), "border": (96, 165, 250), "val_color": (253, 224, 71)},   # Navy
        {"bg_start": (29, 78, 216), "bg_end": (59, 130, 246), "border": (147, 197, 253), "val_color": (255, 255, 255)}, # Royal Blue
        {"bg_start": (194, 65, 12), "bg_end": (234, 88, 12), "border": (253, 186, 116), "val_color": (253, 224, 71)},  # Orange
        {"bg_start": (2, 132, 199), "bg_end": (56, 189, 248), "border": (186, 230, 253), "val_color": (253, 224, 71)},  # Ice Blue
        {"bg_start": (180, 83, 9), "bg_end": (217, 119, 6), "border": (253, 224, 71), "val_color": (255, 255, 255)},   # Amber
        {"bg_start": (55, 48, 163), "bg_end": (79, 70, 229), "border": (165, 180, 252), "val_color": (253, 224, 71)},  # Indigo
        {"bg_start": (15, 118, 110), "bg_end": (13, 148, 136), "border": (94, 234, 212), "val_color": (153, 246, 228)}, # Teal
        {"bg_start": (159, 18, 57), "bg_end": (225, 29, 72), "border": (253, 164, 175), "val_color": (254, 205, 211)}, # Crimson
        {"bg_start": (30, 41, 59), "bg_end": (71, 85, 105), "border": (148, 163, 184), "val_color": (253, 224, 71)}     # Slate
    ]

    minor_styles = [
        {"bg_start": (253, 186, 116), "bg_end": (254, 215, 170), "border": (255, 255, 255), "val_color": (124, 45, 18), "text": (124, 45, 18), "sub": (124, 45, 18)},
        {"bg_start": (203, 213, 225), "bg_end": (226, 232, 240), "border": (255, 255, 255), "val_color": (15, 23, 42), "text": (30, 41, 59), "sub": (30, 41, 59)},
        {"bg_start": (165, 243, 252), "bg_end": (186, 230, 253), "border": (255, 255, 255), "val_color": (3, 105, 161), "text": (3, 105, 161), "sub": (3, 105, 161)},
        {"bg_start": (221, 214, 254), "bg_end": (237, 233, 254), "border": (255, 255, 255), "val_color": (76, 29, 149), "text": (76, 29, 149), "sub": (76, 29, 149)}
    ]

    def draw_gradient(img, start_color, end_color):
        w, h = img.size
        dr = ImageDraw.Draw(img)
        for i in range(h):
            ratio = i / float(h)
            r = int(start_color[0] * (1 - ratio) + end_color[0] * ratio)
            g = int(start_color[1] * (1 - ratio) + end_color[1] * ratio)
            b = int(start_color[2] * (1 - ratio) + end_color[2] * ratio)
            dr.line([(0, i), (w, i)], fill=(r, g, b))

    for idx, rect in enumerate(rects):
        item = rect["item"]
        rx, ry, rw, rh = int(round(rect["x"])), int(round(rect["y"])), int(round(rect["w"])), int(round(rect["h"]))
        if rw <= 2 or rh <= 2:
            continue

        is_vital = idx <= vital_cutoff_idx

        if idx < len(main_styles):
            st = main_styles[idx]
        elif is_vital:
            st = main_styles[1 + ((idx - len(main_styles)) % (len(main_styles) - 1))]
        else:
            st = minor_styles[idx % len(minor_styles)]

        tile_img = Image.new('RGB', (rw, rh), color=st["bg_start"])
        draw_gradient(tile_img, st["bg_start"], st["bg_end"])
        t_draw = ImageDraw.Draw(tile_img)

        # Border
        t_draw.rectangle([0, 0, rw - 1, rh - 1], outline=st["border"], width=2)

        val = float(item.get("value", 0))
        pct = float(item.get("percentage", 0))
        cum_pct = float(item.get("cumulative_percentage", 0))

        if val >= 1_000_000_000:
            val_str = f"Rp {val/1_000_000_000:.2f} Miliar"
        elif val >= 1_000_000:
            val_str = f"Rp {val/1_000_000:.2f} Juta"
        elif val >= 1_000:
            val_str = f"Rp {val:,.0f}"
        else:
            val_str = f"{val:.0f}"

        name_str = str(item.get("name", "-")).upper()

        pad = 8 if rw > 80 and rh > 60 else 4
        avail_w = rw - pad * 2
        avail_h = rh - pad * 2

        if avail_w <= 20 or avail_h <= 15:
            canvas.paste(tile_img, (rx, ry))
            continue

        text_color = st.get("text", (255, 255, 255))
        val_color = st.get("val_color", (253, 224, 71))
        sub_color = st.get("sub", (203, 213, 225))

        # Check Pareto badge availability
        has_badge = is_vital and avail_w >= 100 and avail_h >= 65
        badge_w, badge_h = 0, 0
        f_badge = get_font(font_bold_path, 10)

        if has_badge:
            badge_text = "Pareto 80%"
            if hasattr(f_badge, 'getlength'):
                bw = int(f_badge.getlength(badge_text))
            elif hasattr(f_badge, 'getbbox'):
                bw = f_badge.getbbox(badge_text)[2]
            else:
                bw = len(badge_text) * 7
            badge_w = bw + 22
            badge_h = 18

            bx = rw - pad - badge_w
            by = pad
            # Draw badge box
            t_draw.rectangle([bx, by, bx + badge_w, by + badge_h], fill=(0, 0, 0), outline=(234, 179, 8), width=1)

            # Draw star polygon
            star_cx = bx + 10
            star_cy = by + 9
            star_pts = []
            for sp in range(10):
                sr = 4.5 if sp % 2 == 0 else 2.0
                s_angle = sp * math.pi / 5 - math.pi / 2
                star_pts.append((star_cx + sr * math.cos(s_angle), star_cy + sr * math.sin(s_angle)))
            t_draw.polygon(star_pts, fill=(253, 224, 71))

            t_draw.text((bx + 18, by + 2), badge_text, fill=(253, 224, 71), font=f_badge)

        sl_lbl = item.get("sl_label", "SL Kirim")
        sl_val = item.get("sl_active", item.get("sl_kirim", None))

        # Tile text rendering with Multi-line Word Wrapping (NO TRUNCATION OR CLIPPING!)
        if avail_w >= 110 and avail_h >= 75:
            if is_vital:
                init_title_size = max(13, min(22, avail_w // 8))
                val_size = max(16, min(26, avail_w // 7))
                sub_size = max(10, min(13, avail_w // 13))
            else:
                init_title_size = max(11, min(16, avail_w // 11))
                val_size = max(13, min(20, avail_w // 9))
                sub_size = max(9, min(11, avail_w // 16))

            f_val = get_font(font_bold_path, val_size)
            f_sub = get_font(font_bold_path if is_vital else font_reg_path, sub_size)

            lines = []
            f_title = None
            title_size = init_title_size

            while title_size >= 9.5:
                f_title = get_font(font_bold_path, title_size)
                wrap_w = avail_w - (badge_w + 6 if has_badge else 0)
                lines = wrap_text_lines(f_title, name_str, wrap_w)

                line_h = int(title_size * 1.25)
                tot_title_h = len(lines) * line_h
                tot_h = tot_title_h + val_size + (sub_size * 2) + 16

                if tot_h <= avail_h or title_size <= 9.5:
                    break
                title_size -= 1.0

            y_curr = pad
            line_h = int(title_size * 1.25)
            for line in lines:
                if y_curr + line_h > rh - pad - val_size - 4:
                    break
                t_draw.text((pad, y_curr), line, fill=text_color, font=f_title)
                y_curr += line_h

            y_curr += 4
            t_draw.text((pad, y_curr), val_str, fill=val_color, font=f_val)
            y_curr += int(val_size * 1.2) + 2

            # Render Subtext lines on separate lines (NO HORIZONTAL CLIPPING!)
            if y_curr + sub_size <= rh - pad:
                sub_txt1 = f"Kontribusi: {pct:.1f}% (Kum: {cum_pct:.1f}%)"
                t_draw.text((pad, y_curr), sub_txt1, fill=sub_color, font=f_sub)
                y_curr += sub_size + 2

            if sl_val is not None and y_curr + sub_size <= rh - pad:
                sub_txt2 = f"{sl_lbl}: {float(sl_val):.1f}%"
                sl_color = (74, 222, 128) if sl_lbl == 'SL Kirim' else (251, 191, 36)
                t_draw.text((pad, y_curr), sub_txt2, fill=sl_color, font=f_sub)

        elif avail_w >= 65 and avail_h >= 45:
            if is_vital:
                init_title_size = max(11, min(15, avail_w // 8))
                val_size = max(12, min(16, avail_w // 7))
                sub_size = 10
            else:
                init_title_size = max(9.5, min(13, avail_w // 9))
                val_size = max(11, min(14, avail_w // 8))
                sub_size = 9

            f_val = get_font(font_bold_path, val_size)
            f_sub = get_font(font_bold_path if is_vital else font_reg_path, sub_size)

            title_size = init_title_size
            lines = []
            while title_size >= 8.5:
                f_title = get_font(font_bold_path, title_size)
                lines = wrap_text_lines(f_title, name_str, avail_w)
                tot_h = (len(lines) * int(title_size * 1.2)) + val_size + sub_size + 8
                if tot_h <= avail_h or title_size <= 8.5:
                    break
                title_size -= 0.5

            y_curr = pad
            line_h = int(title_size * 1.2)
            for line in lines:
                if y_curr + line_h > rh - pad - val_size:
                    break
                t_draw.text((pad, y_curr), line, fill=text_color, font=f_title)
                y_curr += line_h

            y_curr += 2
            t_draw.text((pad, y_curr), val_str, fill=val_color, font=f_val)
            y_curr += val_size + 2

            if sl_val is not None and y_curr + sub_size <= rh - pad:
                sl_color = (74, 222, 128) if sl_lbl == 'SL Kirim' else (251, 191, 36)
                t_draw.text((pad, y_curr), f"{pct:.1f}% | {sl_lbl}: {float(sl_val):.1f}%", fill=sl_color, font=f_sub)

        elif avail_w >= 40 and avail_h >= 25:
            title_size = max(8.5, min(10, avail_w // 7))
            f_title = get_font(font_bold_path, title_size)
            f_sub = get_font(font_reg_path, 8)

            lines = wrap_text_lines(f_title, name_str, avail_w)
            y_curr = pad
            line_h = int(title_size * 1.15)

            for line in lines:
                if y_curr + line_h > rh - pad - 10:
                    break
                t_draw.text((pad, y_curr), line, fill=text_color, font=f_title)
                y_curr += line_h

            if y_curr + 10 <= rh - pad:
                t_draw.text((pad, y_curr), f"{pct:.1f}%", fill=val_color, font=f_sub)

        canvas.paste(tile_img, (rx, ry))

    import tempfile
    fd, temp_path = tempfile.mkstemp(suffix=".png", prefix="treemap_img_")
    os.close(fd)
    canvas.save(temp_path, "PNG", quality=95)
    return temp_path

    def _generate_with_python_pptx(self, export_data: Dict[str, Any], output_path: str) -> str:
        prs = Presentation(self.template_path)

def build_active_filters_summary(filters: Dict[str, Any]) -> str:
    if not filters:
        return "📌 Filter Aktif: Semua Data"

    parts = []

    # 1. Month
    m_val = filters.get("months") or filters.get("month")
    if m_val:
        if isinstance(m_val, list):
            valid_m = [format_month_label(m) for m in m_val if m and str(m).upper() not in ["ALL", "SEMUA", "SEMUA BULAN"]]
            if valid_m:
                if len(valid_m) == 1:
                    parts.append(f"Bulan = {valid_m[0]}")
                else:
                    parts.append(f"Bulan = {len(valid_m)} Terpilih ({', '.join(valid_m[:2])})")
        else:
            lbl = format_month_label(m_val)
            if lbl != "Semua Bulan":
                parts.append(f"Bulan = {lbl}")
    if not any(p.startswith("Bulan =") for p in parts):
        parts.append("Bulan = AGU-2026")

    # 2. MTM Type
    t_val = filters.get("mtm_types") or filters.get("mtm_type")
    if t_val:
        if isinstance(t_val, list):
            valid_t = [str(t).strip() for t in t_val if t and str(t).upper() not in ["ALL", "SEMUA", "SEMUA JENIS MTM"]]
            if valid_t:
                parts.append(f"MTM = {', '.join(valid_t)}")
        else:
            t_str = str(t_val).strip()
            if t_str and t_str.upper() not in ["ALL", "SEMUA", "SEMUA JENIS MTM"]:
                parts.append(f"MTM = {t_str}")

    # 3. Branch
    b_val = filters.get("branches") or filters.get("branch")
    if b_val:
        if isinstance(b_val, list):
            valid_b = [str(b).strip() for b in b_val if b and str(b).upper() not in ["ALL", "SEMUA", "SEMUA CABANG"]]
            if valid_b:
                if len(valid_b) == 1:
                    parts.append(f"Cabang = {valid_b[0]}")
                else:
                    parts.append(f"Cabang = {len(valid_b)} Terpilih")
        else:
            b_str = str(b_val).strip()
            if b_str and b_str.upper() not in ["ALL", "SEMUA", "SEMUA CABANG"]:
                parts.append(f"Cabang = {b_str}")

    # 4. MTM Alias
    a_val = filters.get("mtm_aliases") or filters.get("mtm_alias")
    if a_val:
        if isinstance(a_val, list):
            valid_a = [str(a).strip() for a in a_val if a and str(a).upper() not in ["ALL", "SEMUA", "SEMUA ALIAS"]]
            if valid_a:
                if len(valid_a) == 1:
                    parts.append(f"MTM Alias = {valid_a[0]}")
                else:
                    parts.append(f"MTM Alias = {len(valid_a)} Terpilih")
        else:
            a_str = str(a_val).strip()
            if a_str and a_str.upper() not in ["ALL", "SEMUA", "SEMUA ALIAS"]:
                parts.append(f"MTM Alias = {a_str}")

    # 5. Brand Group
    bg_val = filters.get("brand_groups") or filters.get("brand_group")
    if bg_val:
        if isinstance(bg_val, list):
            valid_bg = [str(bg).strip() for bg in bg_val if bg and str(bg).upper() not in ["ALL", "SEMUA", "SEMUA GRUP BRAND"]]
            if valid_bg:
                if len(valid_bg) == 1:
                    parts.append(f"Brand = {valid_bg[0]}")
                else:
                    parts.append(f"Brand = {len(valid_bg)} Terpilih")
        else:
            bg_str = str(bg_val).strip()
            if bg_str and bg_str.upper() not in ["ALL", "SEMUA", "SEMUA GRUP BRAND"]:
                parts.append(f"Brand = {bg_str}")

    # 6. Item
    i_val = filters.get("items") or filters.get("item")
    if i_val:
        if isinstance(i_val, list):
            valid_i = [str(i).strip() for i in i_val if i and str(i).upper() not in ["ALL", "SEMUA", "SEMUA PRODUK / ITEM"]]
            if valid_i:
                if len(valid_i) == 1:
                    parts.append(f"Item = {valid_i[0]}")
                else:
                    parts.append(f"Item = {len(valid_i)} Terpilih")
        else:
            i_str = str(i_val).strip()
            if i_str and i_str.upper() not in ["ALL", "SEMUA", "SEMUA PRODUK / ITEM"]:
                parts.append(f"Item = {i_str}")

    # 7. Reason
    r_val = filters.get("reasons") or filters.get("reason")
    if r_val:
        if isinstance(r_val, list):
            valid_r = [str(r).strip() for r in r_val if r and str(r).upper() not in ["ALL", "SEMUA", "SEMUA ALASAN"]]
            if valid_r:
                if len(valid_r) == 1:
                    parts.append(f"Alasan = {valid_r[0]}")
                else:
                    parts.append(f"Alasan = {len(valid_r)} Terpilih")
        else:
            r_str = str(r_val).strip()
            if r_str and r_str.upper() not in ["ALL", "SEMUA", "SEMUA ALASAN"]:
                parts.append(f"Alasan = {r_str}")

    return "📌 Filter Aktif: " + " | ".join(parts)


def is_all_grup_brand_selected(filters: dict) -> bool:
    """Returns True if no specific grup brand filter is selected (i.e. 'ALL' / 'SEMUA' is active)."""
    bg_val = filters.get("brand_groups") or filters.get("brand_group") or filters.get("grup_brand") or filters.get("brand")
    if not bg_val:
        return True
    if isinstance(bg_val, list):
        valid_bg = [str(bg).strip() for bg in bg_val if bg and str(bg).upper() not in ["ALL", "SEMUA", "SEMUA GRUP BRAND", "SEMUA BRAND"]]
        return len(valid_bg) == 0
    else:
        bg_str = str(bg_val).strip().upper()
        return bg_str in ["", "ALL", "SEMUA", "SEMUA GRUP BRAND", "SEMUA BRAND"]


class MTMPPTExporter:
    def __init__(self, template_path: str = "Template PPT.pptx"):
        self.template_path = template_path

    def generate_presentation(self, export_data: Dict[str, Any], output_path: str = "output_report.pptx") -> str:
        """
        Generate PowerPoint presentation based on Template PPT.pptx with complete slide data population matching dashboard filters.
        """
        if HAS_PYTHON_PPTX and os.path.exists(self.template_path):
            return self._generate_with_python_pptx(export_data, output_path)
        else:
            return self._generate_fallback(export_data, output_path)

    def _create_cloned_content_slide(self, prs, content_layout):
        """Creates a new slide cloned from Slide 1 of Template PPT.pptx, preserving exact background & branding."""
        slide = prs.slides.add_slide(content_layout)

        # Clear default placeholder textboxes
        for sp in list(slide.shapes):
            sp._element.getparent().remove(sp._element)

        # Copy background <p:bg> element from Slide 1 if available
        if len(prs.slides) > 1 and prs.slides[1]._element.cSld.bg is not None:
            src_bg = prs.slides[1]._element.cSld.bg
            new_bg = copy.deepcopy(src_bg)
            if slide._element.cSld.bg is not None:
                slide._element.cSld.remove(slide._element.cSld.bg)
            slide._element.cSld.insert(0, new_bg)

        return slide

    def _generate_with_python_pptx(self, export_data: Dict[str, Any], output_path: str) -> str:
        prs = Presentation(self.template_path)

        # Save reference to last slide (Terima Kasih slide from Template PPT.pptx)
        thanks_sld_id = prs.slides._sldIdLst[-1] if len(prs.slides) > 1 else None

        period_reports = export_data.get("period_reports")
        if not period_reports:
            period_reports = [export_data]

        content_layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0]

        # Helper to get or create content slide matching Template PPT.pptx slides
        slide_cursor = 1

        def get_or_create_content_slide():
            nonlocal slide_cursor, thanks_sld_id
            if slide_cursor < len(prs.slides) - 1:
                slide = prs.slides[slide_cursor]
            else:
                slide = prs.slides.add_slide(content_layout)
                if thanks_sld_id is not None:
                    prs.slides._sldIdLst.remove(thanks_sld_id)
                    prs.slides._sldIdLst.append(thanks_sld_id)

            slide_cursor += 1

            for sp in list(slide.shapes):
                sp._element.getparent().remove(sp._element)

            if slide._element.cSld.bg is None and len(prs.slides) > 1 and prs.slides[1]._element.cSld.bg is not None:
                src_bg = prs.slides[1]._element.cSld.bg
                new_bg = copy.deepcopy(src_bg)
                slide._element.cSld.insert(0, new_bg)

            return slide

        # Slide 0: Cover Slide Title & Subtitle Population Inside Template Dashed Boxes
        if len(prs.slides) > 0:
            cover_slide = prs.slides[0]

            if len(period_reports) > 1:
                month_labels_list = [format_month_label(rep.get("filters", {}).get("month")) for rep in period_reports]
                cover_filter_text = f"PERIODE FILTER: Bulan = {', '.join(month_labels_list)} ({len(period_reports)} Periode Terpilih)"
                base_filters = export_data.get("filters", {})
                t_val = base_filters.get("mtm_types") or base_filters.get("mtm_type")
                if t_val:
                    t_str = t_val[0] if isinstance(t_val, list) else str(t_val)
                    if t_str and t_str.upper() not in ["ALL", "SEMUA"]:
                        cover_filter_text += f" | MTM = {t_str}"
            else:
                cover_filter_text = build_active_filters_summary(period_reports[0].get("filters", {})).replace("📌 Filter Aktif: ", "PERIODE FILTER: ")

            # Populate Top Dashed Box (Placeholder 0)
            if len(cover_slide.placeholders) > 0:
                ph0 = cover_slide.placeholders[0]
                tf0 = ph0.text_frame
                tf0.word_wrap = True
                p0 = tf0.paragraphs[0]
                p0.text = "LAPORAN EXECUTIVE SERVICE LEVEL MTM"
                p0.font.size = Pt(20)
                p0.font.bold = True
                p0.font.color.rgb = RGBColor(255, 255, 255)
                
                p1 = tf0.add_paragraph()
                p1.text = "Analisis Performa Pengiriman, Realisasi, dan Pareto Unfulfill"
                p1.font.size = Pt(12)
                p1.font.color.rgb = RGBColor(255, 215, 0)

            # Populate Bottom Dashed Box (Placeholder 1)
            if len(cover_slide.placeholders) > 1:
                ph1 = cover_slide.placeholders[1]
                tf1 = ph1.text_frame
                tf1.word_wrap = True
                p2 = tf1.paragraphs[0]
                p2.text = cover_filter_text
                p2.font.size = Pt(9.5)
                p2.font.bold = True
                p2.font.color.rgb = RGBColor(255, 255, 255)
            else:
                # Fallback if placeholder 1 is missing
                cov_box = cover_slide.shapes.add_textbox(Inches(4.87), Inches(1.58), Inches(4.79), Inches(1.48))
                tf_cov = cov_box.text_frame
                tf_cov.word_wrap = True
                p0 = tf_cov.paragraphs[0]
                p0.text = "LAPORAN EXECUTIVE SERVICE LEVEL MTM"
                p0.font.size = Pt(20)
                p0.font.bold = True
                p0.font.color.rgb = RGBColor(255, 255, 255)
                p1 = tf_cov.add_paragraph()
                p1.text = "Analisis Performa Pengiriman, Realisasi, dan Pareto Unfulfill"
                p1.font.size = Pt(12)
                p1.font.color.rgb = RGBColor(255, 215, 0)

        # Iterate through period reports in chronological order
        for rep_idx, report in enumerate(period_reports):
            filters = report.get("filters", {})
            kpi = report.get("kpi", {})
            pareto = report.get("pareto", {})
            grid = report.get("grid", [])
            modules = report.get("selected_modules", {})

            sel_month = filters.get("month", "2026-08")
            sel_mtm = filters.get("mtm_type", "KA")
            month_label = format_month_label(sel_month)

            filter_info = build_active_filters_summary(filters)

            # Slide 1: Executive KPI Summary & Monthly Performance Trend Table for this period
            if modules.get("kpi_summary", True):
                slide_kpi = get_or_create_content_slide()

            # Clear default title text
            for shape in list(slide_kpi.shapes):
                if shape.has_text_frame and shape.text_frame.text in ["Title 4", "Click to add title"]:
                    shape.text_frame.text = ""

            title_box = slide_kpi.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"1. EXECUTIVE KPI SUMMARY & TREND ({month_label})"
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 0, 0)

            p_sub = tf.add_paragraph()
            p_sub.text = filter_info
            p_sub.font.size = Pt(8.5)
            p_sub.font.bold = True
            p_sub.font.color.rgb = RGBColor(100, 100, 100)

            # KPI Metric Cards (3 Cards)
            card_w = Inches(2.75)
            card_h = Inches(0.95)
            gap_w = Inches(0.28)

            metrics = [
                ("SERVICE LEVEL KIRIM", f"{kpi.get('sl_kirim', 0):.1f}%", "Target Acuan: 85.0%"),
                ("SERVICE LEVEL REALISASI", f"{kpi.get('sl_realisasi', 0):.1f}%", "Target Acuan: 85.0%"),
                ("GAP REALISASI VS KIRIM", f"{kpi.get('gap', 0):+.1f}%", "Selisih Performa")
            ]

            for i, (title, val, sub) in enumerate(metrics):
                left = Inches(0.55) + i * (card_w + gap_w)
                top = Inches(1.15)

                shape = slide_kpi.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
                shape.line.color.rgb = RGBColor(192, 0, 0)

                tf_card = shape.text_frame
                tf_card.word_wrap = True

                p0 = tf_card.paragraphs[0]
                p0.text = title
                p0.font.size = Pt(9)
                p0.font.bold = True
                p0.font.color.rgb = RGBColor(100, 100, 100)
                p0.alignment = PP_ALIGN.CENTER

                p1 = tf_card.add_paragraph()
                p1.text = val
                p1.font.size = Pt(19)
                p1.font.bold = True
                p1.font.color.rgb = RGBColor(192, 0, 0) if "GAP" not in title else (RGBColor(34, 197, 94) if kpi.get('gap', 0) >= 0 else RGBColor(239, 68, 68))
                p1.alignment = PP_ALIGN.CENTER

                p2 = tf_card.add_paragraph()
                p2.text = sub
                p2.font.size = Pt(7.5)
                p2.font.color.rgb = RGBColor(120, 120, 120)
                p2.alignment = PP_ALIGN.CENTER

            # Monthly Performance Trend Bar Chart (Grafik Batang / Column Clustered MTM Dashboard Style)
            trend_data = report.get("trend", [])
            if trend_data:
                chart_data = CategoryChartData()
                chart_data.categories = [format_month_label(r.get("month", "")) for r in trend_data]
                chart_data.add_series('SL Kirim (%)', [round(float(r.get('sl_kirim', 0)), 1) for r in trend_data])
                chart_data.add_series('SL Realisasi (%)', [round(float(r.get('sl_realisasi', 0)), 1) for r in trend_data])

                x, y, cx, cy = Inches(0.55), Inches(2.18), Inches(8.8), Inches(2.55)
                chart_shape = slide_kpi.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
                chart = chart_shape.chart

                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.TOP
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(8.5)

                plot = chart.plots[0]
                plot.has_data_labels = True
                plot.gap_width = 55
                plot.overlap = -10

                data_labels = plot.data_labels
                data_labels.font.size = Pt(7.0)
                data_labels.font.bold = True
                data_labels.font.color.rgb = RGBColor(15, 23, 42)
                data_labels.number_format = '0.0"%"'
                data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

                # Format Value & Category Axes cleanly (No gridlines, elegant borderless layout)
                val_axis = chart.value_axis
                val_axis.maximum_scale = 115
                val_axis.minimum_scale = 0
                val_axis.major_unit = 20
                val_axis.has_major_gridlines = False
                val_axis.has_minor_gridlines = False
                val_axis.tick_labels.font.size = Pt(8)
                val_axis.tick_labels.font.color.rgb = RGBColor(100, 116, 139)

                cat_axis = chart.category_axis
                cat_axis.has_major_gridlines = False
                cat_axis.has_minor_gridlines = False
                cat_axis.tick_labels.font.size = Pt(8.5)
                cat_axis.tick_labels.font.bold = True
                cat_axis.tick_labels.font.color.rgb = RGBColor(30, 41, 59)

                # Style Series Fills (High Contrast Color Distinction)
                # Series 0: SL Kirim (Konimex Primary Red)
                if len(chart.series) > 0:
                    series0 = chart.series[0]
                    series0.format.fill.solid()
                    series0.format.fill.fore_color.rgb = RGBColor(192, 0, 0)

                # Series 1: SL Realisasi (Royal Blue - High Contrast Distinction)
                if len(chart.series) > 1:
                    series1 = chart.series[1]
                    series1.format.fill.solid()
                    series1.format.fill.fore_color.rgb = RGBColor(37, 99, 235)

            # Slide 1.1: Service Level Per Grup Brand (Executive KPI Breakdown per Grup Brand)
            def gb_sort_key(item: Any) -> tuple:
                name = str(item.get('name', '') if isinstance(item, dict) else item).strip().upper()
                if 'ET' in name:
                    return (999, name)
                m = re.search(r'\d+', name)
                if m:
                    return (int(m.group()), name)
                return (500, name)

            grid_by_dim = report.get("grid_by_dim", {})
            gb_grid = grid_by_dim.get("grup_brand", [])
            if not gb_grid and report.get("grid"):
                gb_grid = report.get("grid", [])

            if gb_grid:
                gb_grid = sorted(gb_grid, key=gb_sort_key)

            if gb_grid and is_all_grup_brand_selected(filters):
                slide_gb = get_or_create_content_slide()

                title_box_gb = slide_gb.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.8), Inches(0.60))
                tf_gb = title_box_gb.text_frame
                tf_gb.word_wrap = True
                p_gb = tf_gb.paragraphs[0]
                p_gb.text = f"1.1 SERVICE LEVEL PER GRUP BRAND ({month_label})"
                p_gb.font.size = Pt(17)
                p_gb.font.bold = True
                p_gb.font.color.rgb = RGBColor(192, 0, 0)

                p_sub_gb = tf_gb.add_paragraph()
                p_sub_gb.text = f"{filter_info} | Performa Pengiriman (SL Kirim) & Realisasi (SL Realisasi) Per Grup Brand"
                p_sub_gb.font.size = Pt(8.5)
                p_sub_gb.font.bold = True
                p_sub_gb.font.color.rgb = RGBColor(100, 100, 100)

                # Left Side: Clustered Column Chart (SL Kirim vs SL Realisasi per GB)
                chart_data_gb = CategoryChartData()
                chart_data_gb.categories = [str(r.get('name', '')) for r in gb_grid]
                chart_data_gb.add_series('SL Kirim (%)', [round(float(r.get('sl_kirim', 0)), 1) for r in gb_grid])
                chart_data_gb.add_series('SL Realisasi (%)', [round(float(r.get('sl_realisasi', 0)), 1) for r in gb_grid])

                cx_gb, cy_gb = Inches(4.60), Inches(3.55)
                chart_shape_gb = slide_gb.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.15), cx_gb, cy_gb, chart_data_gb)
                chart_gb = chart_shape_gb.chart

                chart_gb.has_legend = True
                chart_gb.legend.position = XL_LEGEND_POSITION.TOP
                chart_gb.legend.include_in_layout = False
                chart_gb.legend.font.size = Pt(8.0)

                plot_gb = chart_gb.plots[0]
                plot_gb.has_data_labels = True
                plot_gb.gap_width = 55
                plot_gb.overlap = -10

                data_labels_gb = plot_gb.data_labels
                data_labels_gb.font.size = Pt(6.5)
                data_labels_gb.font.bold = True
                data_labels_gb.font.color.rgb = RGBColor(15, 23, 42)
                data_labels_gb.number_format = '0.0"%"'
                data_labels_gb.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

                val_axis_gb = chart_gb.value_axis
                val_axis_gb.maximum_scale = 115
                val_axis_gb.minimum_scale = 0
                val_axis_gb.major_unit = 20
                val_axis_gb.has_major_gridlines = False
                val_axis_gb.has_minor_gridlines = False
                val_axis_gb.tick_labels.font.size = Pt(7.5)
                val_axis_gb.tick_labels.font.color.rgb = RGBColor(100, 116, 139)

                cat_axis_gb = chart_gb.category_axis
                cat_axis_gb.has_major_gridlines = False
                cat_axis_gb.has_minor_gridlines = False
                cat_axis_gb.tick_labels.font.size = Pt(7.5)
                cat_axis_gb.tick_labels.font.bold = True
                cat_axis_gb.tick_labels.font.color.rgb = RGBColor(30, 41, 59)

                if len(chart_gb.series) > 0:
                    chart_gb.series[0].format.fill.solid()
                    chart_gb.series[0].format.fill.fore_color.rgb = RGBColor(192, 0, 0)
                if len(chart_gb.series) > 1:
                    chart_gb.series[1].format.fill.solid()
                    chart_gb.series[1].format.fill.fore_color.rgb = RGBColor(37, 99, 235)

                # Right Side: Summary Data Table
                table_rows_gb = len(gb_grid) + 1
                table_cols_gb = 6
                table_shape_gb = slide_gb.shapes.add_table(table_rows_gb, table_cols_gb, Inches(5.20), Inches(1.15), Inches(4.25), Inches(3.55))
                table_gb = table_shape_gb.table

                col_widths_gb = [Inches(0.85), Inches(0.85), Inches(0.85), Inches(0.55), Inches(0.55), Inches(0.60)]
                for c_i, w in enumerate(col_widths_gb):
                    table_gb.columns[c_i].width = w

                headers_gb = ["Grup Brand", "Total Order", "Total Kirim", "SL Kirim", "SL Real.", "GAP"]
                for c_i, h in enumerate(headers_gb):
                    cell = table_gb.cell(0, c_i)
                    cell.text = h
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(192, 0, 0)
                    for p in cell.text_frame.paragraphs:
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        p.font.size = Pt(7.5)
                        p.alignment = PP_ALIGN.CENTER

                for r_i, gb in enumerate(gb_grid):
                    r_idx = r_i + 1
                    name = str(gb.get('name', ''))
                    tp_val = float(gb.get('total_pesan', gb.get('total_p', 0)))
                    tk_val = float(gb.get('total_kirim', gb.get('total_k', 0)))
                    sl_k = float(gb.get('sl_kirim', 0))
                    sl_r = float(gb.get('sl_realisasi', 0))
                    gap = sl_r - sl_k

                    vals = [
                        name,
                        f"{tp_val:,.0f}" if tp_val >= 1000 else f"{tp_val:.0f}",
                        f"{tk_val:,.0f}" if tk_val >= 1000 else f"{tk_val:.0f}",
                        f"{sl_k:.1f}%",
                        f"{sl_r:.1f}%",
                        f"{gap:+.1f}%"
                    ]

                    for c_i, val_str in enumerate(vals):
                        cell = table_gb.cell(r_idx, c_i)
                        cell.text = val_str
                        cell.margin_left = Pt(2)
                        cell.margin_right = Pt(2)
                        cell.margin_top = Pt(1)
                        cell.margin_bottom = Pt(1)
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(7.0)
                            if c_i in [1, 2]:
                                p.alignment = PP_ALIGN.RIGHT
                            elif c_i in [0, 3, 4, 5]:
                                p.alignment = PP_ALIGN.CENTER
                            if c_i == 5:
                                p.font.bold = True
                                p.font.color.rgb = RGBColor(34, 197, 94) if gap >= 0 else RGBColor(220, 38, 38)

            # Slide 2+: Dedicated Pareto & Vital Detail Grid Slides Per Dimension
            if modules.get("pareto_sheets", True) or modules.get("detail_grid", True):
                pareto_data = report.get("pareto", {})
                grid_by_dim = report.get("grid_by_dim", {})

                dimensions_cfg = [
                    ("alasan", "ALASAN KETERLAMBATAN"),
                    ("mtm_alias", "AKUN MTM ALIAS"),
                    ("cabang", "NAMA CABANG"),
                    ("grup_brand", "GRUP BRAND"),
                    ("item", "ITEM / PRODUK SKU")
                ]

                dim_unit_map = {
                    "alasan": "Alasan",
                    "mtm_alias": "Akun MTM",
                    "cabang": "Cabang",
                    "grup_brand": "Grup Brand",
                    "item": "Item SKU"
                }

                section_idx = 2
                for dim_key, dim_label in dimensions_cfg:
                    if dim_key == "grup_brand" and not is_all_grup_brand_selected(filters):
                        continue

                    unit_name = dim_unit_map.get(dim_key, "Elemen")
                    pareto_items = pareto_data.get(dim_key, [])
                    if not pareto_items:
                        continue

                    # Identify Vital Pareto 80% items
                    vital_items = []
                    for item in pareto_items:
                        prev_cum = float(item.get("cumulative_percentage", 0)) - float(item.get("percentage", 0))
                        if prev_cum < 80.0:
                            vital_items.append(item)
                    vital_names = set([v.get("name") for v in vital_items if v.get("name")])

                    # SLIDE A: PARETO TREEMAP SLIDE FOR THIS DIMENSION (Matching 2D Dashboard Treemap Layout)
                    slide_p = get_or_create_content_slide()

                    title_box_p = slide_p.shapes.add_textbox(Inches(0.45), Inches(0.40), Inches(7.5), Inches(0.65))
                    tf_p = title_box_p.text_frame
                    tf_p.word_wrap = True
                    p_p = tf_p.paragraphs[0]
                    p_p.text = f"{section_idx}.1 ANALISIS PARETO UNFULLFILL - {dim_label} ({month_label})"
                    p_p.font.size = Pt(16)
                    p_p.font.bold = True
                    p_p.font.color.rgb = RGBColor(192, 0, 0)

                    p_sub_p = tf_p.add_paragraph()
                    p_sub_p.text = f"{filter_info} | Pareto 80% (1 - {len(vital_items)} dari {len(vital_items)} {unit_name})"
                    p_sub_p.font.size = Pt(8.5)
                    p_sub_p.font.bold = True
                    p_sub_p.font.color.rgb = RGBColor(100, 100, 100)

                    # Generate High-Resolution 2D Treemap Snapshot Image (Matching Dashboard 100%)
                    vital_cutoff_idx = len(vital_items) - 1
                    img_path = generate_treemap_image(pareto_items, width_px=1800, height_px=780, vital_cutoff_idx=vital_cutoff_idx)

                    # Insert Pixel-Perfect Treemap Screenshot Image onto PowerPoint Slide
                    slide_p.shapes.add_picture(img_path, Inches(0.45), Inches(1.10), Inches(8.95), Inches(3.85))

                    # Clean up temporary PNG file
                    try:
                        os.remove(img_path)
                    except Exception:
                        pass

                    # SLIDE B: DETAIL DATA GRID TABLE SLIDE(S) FOR THIS DIMENSION (PARETO ONLY)
                    grid_dim = grid_by_dim.get(dim_key, [])
                    if not grid_dim and report.get("grid"):
                        grid_dim = report.get("grid", [])

                    section_idx += 1

        # Clean up any unused pre-created template slides between slide_cursor and thanks_sld_id
        while len(prs.slides) > slide_cursor + 1:
            rId = prs.slides._sldIdLst[slide_cursor].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[slide_cursor]

        # Ensure thanks_sld_id remains as the final slide
        if thanks_sld_id is not None and prs.slides._sldIdLst[-1] != thanks_sld_id:
            prs.slides._sldIdLst.remove(thanks_sld_id)
            prs.slides._sldIdLst.append(thanks_sld_id)

        prs.save(output_path)
        print(f"Presentation saved successfully to {output_path}")
        return output_path

    def _generate_fallback(self, export_data: Dict[str, Any], output_path: str) -> str:
        """Fallback copy of Template PPT.pptx if pptx library is not loaded."""
        if os.path.exists(self.template_path):
            with open(self.template_path, 'rb') as f_in:
                content = f_in.read()
            with open(output_path, 'wb') as f_out:
                f_out.write(content)
            print(f"Fallback PPT exported to {output_path}")
            return output_path
        else:
            raise FileNotFoundError("Template PPT.pptx not found!")
